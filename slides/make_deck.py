"""Build the project deck: A Layered Agent Fund.

    .venv/bin/pip install python-pptx
    .venv/bin/python slides/make_deck.py

Output: slides/layered_agent_fund_phase1.pptx

Structure: 14 story slides + 5 appendix slides. The story carries three beats —
the layers idea, the model vs the formula, and what the analyst reads. All the
machinery (contracts, DV01, the full diagnostic tables) lives in the appendix,
so the story stays a story and the Q&A stays armed.

Numbers come from reports/phase1_{vector,text,textvec}.md, phase1_7d_textvec.md
and the matching *.audit.json. NOT from phase1.md (superseded).

Logos: drop berkeley.png / blackrock.png into slides/assets/ and re-run.
Theme: Berkeley / Haas — Berkeley Blue #003262 + California Gold #FDB515.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── theme ───────────────────────────────────────────────────────────────────
BERK_BLUE = RGBColor(0x00, 0x32, 0x62)
GOLD = RGBColor(0xFD, 0xB5, 0x15)
FOUNDERS = RGBColor(0x3B, 0x7E, 0xA1)
BAY_FOG = RGBColor(0xDD, 0xD5, 0xC7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6E, 0x6E, 0x6E)
PALE = RGBColor(0xA8, 0xAF, 0xB8)
LIGHT = RGBColor(0xF3, 0xF5, 0xF8)
RED = RGBColor(0x9B, 0x22, 0x26)
GREEN = RGBColor(0x1E, 0x6B, 0x3A)

ARM_DET = RGBColor(0x46, 0x53, 0x5E)     # the formula
ARM_VEC = RGBColor(0x3B, 0x7E, 0xA1)     # numbers
ARM_TXT = RGBColor(0x0E, 0x7C, 0x6B)     # words
ARM_BOTH = RGBColor(0xC4, 0x82, 0x0E)    # both

BODY = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.6)
CONTENT_W = W - 2 * M

HERE = Path(__file__).parent
ASSETS = HERE / "assets"
OUT = HERE / "layered_agent_fund_phase1.pptx"

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
_page = 0


# ── primitives ──────────────────────────────────────────────────────────────
def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         dash=False):
    s = slide.shapes.add_shape(shape, int(x), int(y), int(w), int(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
        if dash:
            s.line.dash_style = 4
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def text(slide, x, y, w, h, runs, size=12, color=INK, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, spacing=1.1, anchor=MSO_ANCHOR.TOP, space_after=4):
    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = [runs] if isinstance(runs, str) else runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after)
        for t, ov in ([(para, {})] if isinstance(para, str) else para):
            r = p.add_run()
            r.text = t
            f = r.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.italic = ov.get("italic", False)
            f.color.rgb = ov.get("color", color)
    return box


def shape_text(s, label, size, color, bold=True, font=BODY):
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return s


def pill(slide, x, y, w, h, label, fill=GOLD, color=BERK_BLUE, size=9, bold=True,
         line=None, dash=False, font=BODY):
    s = rect(slide, x, y, w, h, fill=fill, line=line, dash=dash,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    s.adjustments[0] = 0.3
    return shape_text(s, label, size, color, bold, font)


def arrow(slide, x, y, w, h, fill=BERK_BLUE, shape=MSO_SHAPE.RIGHT_ARROW):
    return rect(slide, x, y, w, h, fill=fill, shape=shape)


def table(slide_, x, y, w, data, col_w, aligns=None, size=10.5, row_h=Inches(0.31),
          fonts=None, cell_color=None, header_fill=BERK_BLUE):
    rows, cols = len(data), len(data[0])
    shape = slide_.shapes.add_table(rows, cols, int(x), int(y), int(w), int(row_h * rows))
    tbl = shape.table
    tbl.first_row = True
    tbl.horz_banding = False
    for i, frac in enumerate(col_w):
        tbl.columns[i].width = int(w * frac)
    for r in range(rows):
        tbl.rows[r].height = int(row_h)
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else (WHITE if r % 2 else LIGHT)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (PP_ALIGN.LEFT if aligns is None else aligns[c])
            run = p.add_run()
            run.text = str(data[r][c])
            f = run.font
            f.size = Pt(size)
            f.name = (fonts[c] if fonts else BODY) if r else BODY
            f.bold = r == 0
            f.color.rgb = WHITE if r == 0 else (cell_color or {}).get((r, c), INK)
    return tbl


L, C, R = PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT


# ── logos + furniture ───────────────────────────────────────────────────────
def _berkeley_wordmark(slide, x, y, h, dark):
    return text(slide, x, y - Inches(0.03), Inches(2.4), h + Inches(0.06),
                [[("Berkeley", {"size": 15, "bold": True, "color": WHITE if dark else BERK_BLUE}),
                  ("  |  ", {"size": 13, "color": GOLD}),
                  ("MFE", {"size": 11, "bold": True, "color": GOLD if dark else FOUNDERS})]],
                anchor=MSO_ANCHOR.MIDDLE)


def _blackrock_wordmark(slide, x, y, h, dark):
    return text(slide, x, y - Inches(0.03), Inches(2.0), h + Inches(0.06),
                [[("BLACKROCK", {"size": 13, "bold": True, "color": WHITE if dark else INK})]],
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def _logo(slide, name, x, y, h, fallback, dark):
    for ext in (".png", ".jpg", ".jpeg"):
        p = ASSETS / f"{name}{ext}"
        if p.exists():
            return slide.shapes.add_picture(str(p), int(x), int(y), height=int(h))
    return fallback(slide, x, y, h, dark)


def logos(slide, y, h=Inches(0.30), dark=False):
    _logo(slide, "berkeley", M, y, h, _berkeley_wordmark, dark)
    _logo(slide, "blackrock", W - M - Inches(2.0), y, h, _blackrock_wordmark, dark)


def slide(title, kicker=None, appendix=False):
    """A story slide: the headline IS the claim. Kicker is optional and quiet."""
    global _page
    _page += 1
    s = blank()
    y = Inches(0.42)
    if kicker:
        text(s, M, Inches(0.34), CONTENT_W, Inches(0.24), kicker.upper(),
             size=10.5, color=GOLD if appendix else FOUNDERS, bold=True)
        y = Inches(0.66)
    text(s, M, y, CONTENT_W, Inches(0.6), title, size=27, color=BERK_BLUE, bold=True,
         spacing=1.05)
    rect(s, M, Inches(1.34), Inches(1.5), Inches(0.055), fill=GOLD)
    rect(s, M, Inches(6.94), CONTENT_W, Emu(9525), fill=BAY_FOG)
    logos(s, Inches(7.05))
    text(s, W / 2 - Inches(0.4), Inches(7.05), Inches(0.8), Inches(0.3), f"{_page}",
         size=10, color=GREY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def punch(s, y, body, color=BERK_BLUE, size=15, h=Inches(0.5)):
    """One plain sentence, the takeaway. No box, no bar — just the line."""
    return text(s, M, y, CONTENT_W, h, body, size=size, color=color, spacing=1.25)


def note(slide_, y, body, fill=LIGHT, bar=GOLD, h=Inches(0.6), w=None):
    w = CONTENT_W if w is None else w
    rect(slide_, M, y, w, h, fill=fill)
    rect(slide_, M, y, Inches(0.055), h, fill=bar)
    text(slide_, M + Inches(0.22), y + Inches(0.09), w - Inches(0.4), h - Inches(0.18),
         body, size=11.5, color=INK, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = blank()
rect(s, 0, 0, W, H, fill=BERK_BLUE)
rect(s, 0, 0, Inches(0.14), H, fill=GOLD)
text(s, Inches(1.1), Inches(2.0), Inches(11), Inches(0.3),
     "UC BERKELEY MFE  ×  BLACKROCK  ·  INDUSTRY PROJECT", size=13, color=GOLD, bold=True)
text(s, Inches(1.1), Inches(2.5), Inches(11), Inches(1.0),
     "A Layered Agent Fund", size=54, color=WHITE, bold=True)
rect(s, Inches(1.1), Inches(3.6), Inches(2.2), Inches(0.05), fill=GOLD)
text(s, Inches(1.1), Inches(3.95), Inches(10.6), Inches(0.9),
     [[("Three layers, two ways to build an analyst, and one surprise about what "
        "actually matters.", {"color": WHITE})]], size=19, spacing=1.3)
text(s, Inches(1.1), Inches(5.75), Inches(6), Inches(0.5),
     [[("Elias Roubache", {"size": 14, "bold": True, "color": WHITE})],
      [("Macro / FICC analyst pool  ·  July 2026", {"size": 12, "color": BAY_FOG})]],
     spacing=1.25)
rect(s, Inches(1.1), Inches(6.55), Inches(11.2), Emu(9525), fill=FOUNDERS)
logos(s, Inches(6.8), h=Inches(0.34), dark=True)

# ════════════════════════════════════════════════════════════════════════════
# 2 — The problem
# ════════════════════════════════════════════════════════════════════════════
s = slide("When an AI fund loses money, nobody can say why")

text(s, M, Inches(1.65), Inches(5.6), Inches(0.25),
     "THE OBVIOUS WAY TO BUILD IT", size=10, color=GREY, bold=True)
for i, name in enumerate(["Warren Buffett", "Ray Dalio", "Aswath Damodaran"]):
    y = Inches(2.35) + i * Inches(0.62)
    pill(s, M, y, Inches(1.85), Inches(0.42), name, fill=WHITE, color=GREY, size=10.5,
         line=PALE)
    arrow(s, M + Inches(1.95), y + Inches(0.11), Inches(0.5), Inches(0.2), fill=PALE)
blob = rect(s, M + Inches(2.6), Inches(2.3), Inches(2.5), Inches(1.9), fill=PALE,
            shape=MSO_SHAPE.OVAL)
shape_text(blob, "one blended\nopinion", 14, WHITE)
text(s, M, Inches(4.55), Inches(5.4), Inches(0.6),
     "Each agent is a whole investor with a whole opinion. A weight blends them "
     "into one book.", size=12.5, color=GREY, spacing=1.25)

rect(s, M + Inches(6.3), Inches(1.6), Emu(9525), Inches(3.6), fill=BAY_FOG)

text(s, M + Inches(6.9), Inches(1.65), Inches(5.2), Inches(0.25),
     "THE MONTH IT LOSES 4%", size=10, color=RED, bold=True)
text(s, M + Inches(6.9), Inches(2.0), Inches(5.2), Inches(0.5),
     "Which part was wrong?", size=24, color=BERK_BLUE, bold=True)
for i, (what, q) in enumerate([
    ("The belief", "was inflation really rising?"),
    ("The trade", "or was it the wrong way to bet on it?"),
    ("The size", "or was the bet simply too big?"),
]):
    y = Inches(2.75) + i * Inches(0.66)
    rect(s, M + Inches(6.9), y, Inches(5.2), Inches(0.52), fill=LIGHT)
    text(s, M + Inches(7.1), y + Inches(0.1), Inches(4.9), Inches(0.32),
         [[(what, {"size": 12.5, "bold": True, "color": INK}),
           ("   " + q, {"size": 12, "color": GREY})]], anchor=MSO_ANCHOR.MIDDLE)
text(s, M + Inches(6.9), Inches(4.8), Inches(5.2), Inches(0.4),
     "You cannot know. Nothing in the machine ever recorded the difference.",
     size=12.5, color=RED, bold=True, spacing=1.25)

punch(s, Inches(5.6),
      [[("The opinions arrive pre-mixed. ", {"bold": True}),
        ("There is no seam to cut along — so there is no question you can ask, and "
         "nothing to fix on purpose.", {})]])

# ════════════════════════════════════════════════════════════════════════════
# 3 — The layers idea  ★
# ════════════════════════════════════════════════════════════════════════════
s = slide("So we gave the fund an org chart")

punch(s, Inches(1.6),
      "A real fund doesn't blend opinions. It holds a meeting — where different people "
      "have different jobs, and each one is answerable for their own.", size=14,
      color=GREY)

BY, BH = Inches(2.35), Inches(2.5)
bw, bgap = Inches(3.5), Inches(0.9)
xs = [M + i * (bw + bgap) for i in range(3)]
layers = [
    ("BELIEVE", "The analysts", ARM_VEC,
     "Each one watches a single force in the world — inflation, jobs, the Fed's "
     "balance sheet — and says what it will do next. Nothing else."),
    ("ACT", "The PM", GOLD,
     "Listens to all of them and turns their beliefs into one trade: the one that "
     "captures what they collectively think and hedges away everything else."),
    ("SURVIVE", "The fund", FOUNDERS,
     "Decides how much risk the trade gets — by how sure it is, how risky it is, and "
     "how much it overlaps everything else we hold."),
]
for x, (verb, who, accent, body) in zip(xs, layers):
    rect(s, x, BY, bw, BH, fill=WHITE, line=BAY_FOG, line_w=1.25)
    rect(s, x, BY, bw, Inches(0.07), fill=accent)
    text(s, x + Inches(0.25), BY + Inches(0.28), bw - Inches(0.5), Inches(0.45), verb,
         size=26, color=BERK_BLUE, bold=True)
    text(s, x + Inches(0.25), BY + Inches(0.82), bw - Inches(0.5), Inches(0.24), who,
         size=12, color=accent, bold=True)
    text(s, x + Inches(0.25), BY + Inches(1.16), bw - Inches(0.5), Inches(1.2), body,
         size=11.5, color=INK, spacing=1.25)

for i, label in enumerate(["a view", "a trade"]):
    gx = xs[i] + bw
    arrow(s, gx + Inches(0.1), BY + Inches(1.1), bgap - Inches(0.2), Inches(0.26))
    text(s, gx, BY + Inches(0.72), bgap, Inches(0.24), label, size=11, color=BERK_BLUE,
         bold=True, align=PP_ALIGN.CENTER)

fb = rect(s, xs[0] + Inches(0.6), Inches(5.05), xs[2] + bw - xs[0] - Inches(1.2),
          Inches(0.36), fill=FOUNDERS, shape=MSO_SHAPE.LEFT_ARROW)
shape_text(fb, "and the capital flows back down", 10.5, WHITE)

text(s, M, Inches(5.62), CONTENT_W, Inches(0.55),
     [[("Today the analysts are macro and rates", {"size": 12, "bold": True,
                                                   "color": BERK_BLUE}),
       ("   inflation · jobs · balance sheet · term premium · curve · breakevens · "
        "financial conditions          ", {"size": 12, "color": GREY}),
       ("equity analysts — teammates, same slot", {"size": 12, "italic": True,
                                                   "color": PALE})]])

punch(s, Inches(6.25),
      [[("Each layer talks to the next through one fixed handoff, and never reaches "
         "inside the others. ", {}),
        ("That is the whole discipline.", {"bold": True})]], size=13)

# ════════════════════════════════════════════════════════════════════════════
# 4 — The boundary
# ════════════════════════════════════════════════════════════════════════════
s = slide("An analyst may say “inflation is rising.” It may not say “buy the 10-year.”")

half = Inches(5.85)
rect(s, M, Inches(1.75), half, Inches(2.9), fill=RGBColor(0xEC, 0xF5, 0xEF))
rect(s, M, Inches(1.75), half, Inches(0.07), fill=GREEN)
text(s, M + Inches(0.3), Inches(2.0), half - Inches(0.6), Inches(0.3),
     "WHAT AN ANALYST IS ALLOWED TO SAY", size=10.5, color=GREEN, bold=True)
text(s, M + Inches(0.3), Inches(2.5), half - Inches(0.6), Inches(1.4),
     [[("“Inflation is rising.", {"size": 21, "bold": True, "color": BERK_BLUE})],
      [("  I'm 46% sure.", {"size": 21, "bold": True, "color": BERK_BLUE})],
      [("  It should hold for about a quarter.”", {"size": 21, "bold": True,
                                                   "color": BERK_BLUE})]], spacing=1.2)
text(s, M + Inches(0.3), Inches(4.05), half - Inches(0.6), Inches(0.4),
     "A claim about the world. It can be checked later against what the world did.",
     size=12, color=GREEN, spacing=1.2)

xr = M + half + Inches(0.45)
rect(s, xr, Inches(1.75), half, Inches(2.9), fill=RGBColor(0xFD, 0xF3, 0xF3))
rect(s, xr, Inches(1.75), half, Inches(0.07), fill=RED)
text(s, xr + Inches(0.3), Inches(2.0), half - Inches(0.6), Inches(0.3),
     "WHAT IT MAY NEVER SAY", size=10.5, color=RED, bold=True)
text(s, xr + Inches(0.3), Inches(2.5), half - Inches(0.6), Inches(1.4),
     [[("“Buy the 10-year.”", {"size": 21, "bold": True, "color": PALE})],
      [("“Put on a flattener.”", {"size": 21, "bold": True, "color": PALE})],
      [("“Go short duration.”", {"size": 21, "bold": True, "color": PALE})]], spacing=1.2)
text(s, xr + Inches(0.3), Inches(4.05), half - Inches(0.6), Inches(0.4),
     "It doesn't know what a flattener is. It has never seen a price. That is on purpose.",
     size=12, color=RED, spacing=1.2)

div = rect(s, M + half + Inches(0.06), Inches(1.95), Inches(0.32), Inches(2.5), fill=GOLD)
text(s, M + half - Inches(1.2), Inches(4.72), Inches(2.8), Inches(0.28),
     "only the PM crosses this line", size=11, color=BERK_BLUE, bold=True,
     align=PP_ALIGN.CENTER)

punch(s, Inches(5.35),
      [[("This is what buys us the answer to slide 2. ", {"bold": True}),
        ("Because the analyst only ever makes a claim about the world, we can grade it "
         "on the world — did inflation actually rise? — completely separately from "
         "whether the trade made money. When the fund is down, we can ask which layer "
         "failed, and get an answer.", {})]], size=14, h=Inches(1.3))

# ════════════════════════════════════════════════════════════════════════════
# 5 — Meet one analyst  (machinery allowed here)
# ════════════════════════════════════════════════════════════════════════════
s = slide("Meet one analyst", kicker="the machinery, once")

punch(s, Inches(1.6), "One driver in, one claim out. Every analyst in the pool is this "
      "shape.", size=14, color=GREY)

FY = Inches(2.35)
# input 1
rect(s, M, FY, Inches(2.1), Inches(1.0), fill=WHITE, line=ARM_VEC, line_w=1.5)
text(s, M + Inches(0.15), FY + Inches(0.14), Inches(1.8), Inches(0.7),
     [[("CPI", {"size": 17, "bold": True, "color": ARM_VEC})],
      [("monthly, from FRED", {"size": 10, "color": GREY})],
      [("the only number it sees", {"size": 10, "italic": True, "color": GREY})]],
     spacing=1.15, space_after=1)
# input 2
rect(s, M, FY + Inches(1.25), Inches(2.1), Inches(1.0), fill=WHITE, line=PALE, line_w=1.25,
     dash=True)
text(s, M + Inches(0.15), FY + Inches(1.39), Inches(1.8), Inches(0.7),
     [[("Fed statement", {"size": 15, "bold": True, "color": PALE})],
      [("the words, not the number", {"size": 10, "color": PALE})],
      [("only when we let it →  slide 9", {"size": 10, "italic": True, "color": ARM_TXT})]],
     spacing=1.15, space_after=1)

arrow(s, M + Inches(2.2), FY + Inches(0.37), Inches(0.5), Inches(0.26), fill=ARM_VEC)
arrow(s, M + Inches(2.2), FY + Inches(1.62), Inches(0.5), Inches(0.26), fill=PALE)

# the gate
gate = rect(s, M + Inches(2.85), FY, Inches(1.55), Inches(2.25), fill=BERK_BLUE)
text(s, M + Inches(2.95), FY + Inches(0.3), Inches(1.35), Inches(1.6),
     [[("🔒", {"size": 22, "color": GOLD})],
      [("the gate", {"size": 13, "bold": True, "color": WHITE})],
      [("It sees only its own data, and never anything dated after the day of the "
        "meeting.", {"size": 9.5, "color": BAY_FOG})]], spacing=1.2, align=PP_ALIGN.CENTER)

arrow(s, M + Inches(4.5), FY + Inches(0.99), Inches(0.5), Inches(0.26))

# the analyst
rect(s, M + Inches(5.15), FY, Inches(2.8), Inches(2.25), fill=LIGHT)
rect(s, M + Inches(5.15), FY, Inches(2.8), Inches(0.07), fill=GOLD)
text(s, M + Inches(5.35), FY + Inches(0.25), Inches(2.4), Inches(1.9),
     [[("The inflation analyst", {"size": 14, "bold": True, "color": BERK_BLUE})],
      [("It measures how fast prices are moving, and how far that is from the Fed's 2% "
        "target.", {"size": 11, "color": INK})],
      [("Then, optionally, it asks Claude to improve on that reading.",
        {"size": 11, "color": INK})],
      [("It has no view on anything else, and it is not asked to agree with anyone.",
        {"size": 10.5, "italic": True, "color": GREY})]], spacing=1.2, space_after=6)

arrow(s, M + Inches(8.05), FY + Inches(0.99), Inches(0.5), Inches(0.26))

# the view
rect(s, M + Inches(8.7), FY, Inches(3.43), Inches(2.25), fill=WHITE, line=BERK_BLUE,
     line_w=1.5)
text(s, M + Inches(8.95), FY + Inches(0.22), Inches(2.95), Inches(1.9),
     [[("Its one claim", {"size": 10.5, "bold": True, "color": FOUNDERS})],
      [("inflation: up", {"size": 19, "bold": True, "color": BERK_BLUE})],
      [("how sure — 0.46 out of 1", {"size": 12, "color": INK})],
      [("for how long — about a quarter", {"size": 12, "color": INK})],
      [("because — “prices are up 3.1% on the year, and still accelerating”",
        {"size": 10.5, "italic": True, "color": GREY})]], spacing=1.2, space_after=5)

note(s, Inches(4.95),
     [[("The two inputs do not arrive the same way. ", {"bold": True}),
       ("The number comes through the gate, and we check on every single run that it "
        "obeyed. The Fed's words come in by their own publication date — we test that "
        "it never gets tomorrow's statement, but that test runs offline, not on the "
        "live run. Worth knowing before slide 11.", {})]], h=Inches(0.72))

punch(s, Inches(5.95),
      [[("Seven analysts, one force each. ", {"bold": True}),
        ("None of them knows what a trade is.", {})]], size=13.5)

# ════════════════════════════════════════════════════════════════════════════
# 6 — The PM
# ════════════════════════════════════════════════════════════════════════════
s = slide("The PM knows one thing the analysts don't")

punch(s, Inches(1.6),
      "Four analysts say four things about the world. None of them names a trade. The "
      "PM's job — and its whole edge — is knowing how those forces reach the market.",
      size=14, color=GREY)

for i, (drv, said) in enumerate([
    ("Inflation", "still hot"), ("Jobs", "still tight"),
    ("Balance sheet", "draining"), ("Term premium", "rising"),
]):
    y = Inches(2.5) + i * Inches(0.6)
    rect(s, M, y, Inches(2.7), Inches(0.48), fill=LIGHT)
    text(s, M + Inches(0.18), y + Inches(0.09), Inches(2.4), Inches(0.3),
         [[(drv, {"size": 11.5, "bold": True, "color": BERK_BLUE}),
           ("  " + said, {"size": 11.5, "color": GREY})]], anchor=MSO_ANCHOR.MIDDLE)

arrow(s, M + Inches(2.85), Inches(3.4), Inches(0.55), Inches(0.28))

rect(s, M + Inches(3.6), Inches(2.4), Inches(3.9), Inches(2.9), fill=WHITE, line=BAY_FOG)
text(s, M + Inches(3.8), Inches(2.55), Inches(3.5), Inches(0.5),
     "So policy stays tighter for longer.", size=13, color=BERK_BLUE, bold=True,
     spacing=1.2)
# the curve picture
base_y = Inches(4.45)
rect(s, M + Inches(3.95), base_y, Inches(3.2), Emu(9525), fill=PALE)
for lbl, cx, ah, col in [("2-year", Inches(4.7), Inches(1.15), RED),
                         ("10-year", Inches(6.5), Inches(0.5), FOUNDERS)]:
    a = rect(s, M + cx, base_y - ah, Inches(0.42), ah, fill=col, shape=MSO_SHAPE.UP_ARROW)
    text(s, M + cx - Inches(0.25), base_y + Inches(0.08), Inches(0.92), Inches(0.22), lbl,
         size=10.5, color=INK, bold=True, align=PP_ALIGN.CENTER)
text(s, M + Inches(3.8), Inches(3.05), Inches(0.9), Inches(0.5),
     "yields\npushed up", size=9.5, color=GREY, spacing=1.15)
text(s, M + Inches(3.8), Inches(4.85), Inches(3.5), Inches(0.4),
     "The front of the curve takes the hit. The long end moves less.",
     size=11.5, color=INK, spacing=1.2)

arrow(s, M + Inches(7.65), Inches(3.4), Inches(0.55), Inches(0.28))

rect(s, M + Inches(8.4), Inches(2.4), Inches(3.73), Inches(2.9), fill=BERK_BLUE)
text(s, M + Inches(8.65), Inches(2.6), Inches(3.25), Inches(2.0),
     [[("One trade", {"size": 11, "bold": True, "color": GOLD})],
      [("A 2s10s flattener", {"size": 21, "bold": True, "color": WHITE})],
      [("Not “sell bonds”. Buy the 10-year, sell the 2-year, sized so that a general "
        "move in rates cancels out.", {"size": 11.5, "color": BAY_FOG})],
      [("Only the thing they actually believe — the gap between the two ends — pays or "
        "loses.", {"size": 11.5, "color": WHITE})]], spacing=1.25, space_after=7)

punch(s, Inches(5.55),
      [[("This map is the PM's edge, and it is the part the analysts cannot hold — none "
         "of them can see more than one force. ", {}),
        ("And when they disagree with each other, the PM bets smaller.",
         {"bold": True})]], size=13.5, h=Inches(0.8))

# ════════════════════════════════════════════════════════════════════════════
# 7 — We built every analyst twice
# ════════════════════════════════════════════════════════════════════════════
s = slide("We built every analyst twice")

punch(s, Inches(1.6),
      "The interesting question isn't whether an AI fund can be built. It's whether the "
      "AI is doing any of the work. So each analyst exists in two versions.", size=14,
      color=GREY)

pw = Inches(5.85)
for i, (name, sub, accent, lines, cost) in enumerate([
    ("THE FORMULA", "no AI at all", ARM_DET,
     ["Read CPI. Measure how much it moved. Say “up”, “down” or “flat”, and how sure "
      "you are. About ten lines of arithmetic.",
      "It cannot know the future — no memory, no training data, nothing to recall. "
      "Whatever it gets right, it got right honestly."],
     "free  ·  instant"),
    ("THE MODEL", "same job, given to Claude", ARM_VEC,
     ["Hand it the exact same reading and ask for a better view. Same mandate, same "
      "output, graded exactly the same way.",
      "It has read the whole internet. On paper, the more sophisticated analyst by a "
      "distance."],
     "$0.94  ·  1,252 calls  ·  47 min"),
]):
    x = M + i * (pw + Inches(0.43))
    rect(s, x, Inches(2.35), pw, Inches(2.75), fill=LIGHT)
    rect(s, x, Inches(2.35), pw, Inches(0.07), fill=accent)
    text(s, x + Inches(0.3), Inches(2.6), pw - Inches(0.6), Inches(0.34), name,
         size=19, color=BERK_BLUE, bold=True)
    text(s, x + Inches(0.3), Inches(3.0), pw - Inches(0.6), Inches(0.24), sub,
         size=12, color=accent, bold=True)
    text(s, x + Inches(0.3), Inches(3.4), pw - Inches(0.6), Inches(1.3), lines,
         size=12, color=INK, spacing=1.25, space_after=7)
    text(s, x + Inches(0.3), Inches(4.72), pw - Inches(0.6), Inches(0.26), cost,
         size=12, color=accent, bold=True, font=MONO)

punch(s, Inches(5.45),
      [[("Then we asked the only question that matters: ", {}),
        ("does the model earn its keep?", {"bold": True, "color": BERK_BLUE})]], size=17)

note(s, Inches(6.1),
     [[("How we grade either of them: ", {"bold": True}),
       ("we ask whether the driver actually moved the way the analyst said it would, "
        "over the next quarter — and we compare that to flipping a coin. 313 weekly "
        "meetings, real history, 2019 to 2024.", {})]], h=Inches(0.55))

# ════════════════════════════════════════════════════════════════════════════
# 8 — The model added nothing  ★ the money slide
# ════════════════════════════════════════════════════════════════════════════
s = slide("The model added nothing")

text(s, M, Inches(1.75), CONTENT_W, Inches(0.35),
     "HOW MUCH BETTER THAN A COIN FLIP, ACROSS 313 MEETINGS", size=11.5, color=GREY,
     bold=True, align=PP_ALIGN.CENTER)

for i, (val, who, sub, accent) in enumerate([
    ("+0.0645", "the formula", "ten lines of arithmetic", ARM_DET),
    ("+0.0648", "the model", "$0.94 · 1,252 calls to Claude", ARM_VEC),
]):
    x = M + Inches(0.4) + i * Inches(6.4)
    text(s, x, Inches(2.35), Inches(5.0), Inches(1.3), val, size=88, color=accent,
         bold=True, align=PP_ALIGN.CENTER, font=MONO)
    text(s, x, Inches(3.75), Inches(5.0), Inches(0.32), who, size=17, color=BERK_BLUE,
         bold=True, align=PP_ALIGN.CENTER)
    text(s, x, Inches(4.12), Inches(5.0), Inches(0.28), sub, size=12, color=GREY,
         align=PP_ALIGN.CENTER)

text(s, W / 2 - Inches(0.9), Inches(2.75), Inches(1.8), Inches(0.4), "vs", size=20,
     color=PALE, bold=True, align=PP_ALIGN.CENTER)
rect(s, W / 2 - Inches(0.75), Inches(3.35), Inches(1.5), Inches(0.32), fill=GOLD)
text(s, W / 2 - Inches(0.75), Inches(3.35), Inches(1.5), Inches(0.32),
     "+0.0003", size=12, color=BERK_BLUE, bold=True, align=PP_ALIGN.CENTER,
     anchor=MSO_ANCHOR.MIDDLE, font=MONO)

punch(s, Inches(4.75),
      [[("Given the same numbers to look at, the model just restated what the formula "
         "already said.", {"size": 17, "bold": True, "color": BERK_BLUE})],
       [("It didn't do worse. It simply didn't do better — and it disagreed with the "
         "formula on only 84 calls out of 1,252.", {"size": 14, "color": INK})]],
      size=14, h=Inches(1.0))

note(s, Inches(6.05),
     [[("This is the result that set the direction. ", {"bold": True}),
       ("If the expensive part of the analyst isn't where the edge comes from, then the "
        "lever is somewhere else — and the only other thing an analyst has is what it "
        "reads.", {})]], h=Inches(0.55))

# ════════════════════════════════════════════════════════════════════════════
# 9 — So we changed what it reads
# ════════════════════════════════════════════════════════════════════════════
s = slide("So we changed what it reads, not what it is")

punch(s, Inches(1.6),
      [[("A CPI print is the past. ", {"bold": True, "color": BERK_BLUE}),
        ("The Fed says what it is going to do — in words, out loud — before it ever "
         "shows up in the data. So we gave the same model three different things to "
         "read, and changed nothing else.", {"color": GREY})]], size=14, h=Inches(0.7))

arms = [
    ("NUMBERS", "what it read before", ARM_VEC,
     "The measurement: prices are up 3.1% on the year, and accelerating.", "$0.94"),
    ("WORDS", "the Fed's own statement", ARM_TXT,
     "The latest FOMC statement, and nothing else. No numbers at all.", "$1.57"),
    ("BOTH", "the number and the story", ARM_BOTH,
     "The measurement and the statement, together.", "$1.75"),
]
aw = Inches(3.85)
for i, (name, sub, accent, body, cost) in enumerate(arms):
    x = M + i * (aw + Inches(0.29))
    rect(s, x, Inches(2.5), aw, Inches(1.95), fill=LIGHT)
    rect(s, x, Inches(2.5), aw, Inches(0.07), fill=accent)
    text(s, x + Inches(0.25), Inches(2.75), aw - Inches(0.5), Inches(0.34), name,
         size=21, color=accent, bold=True)
    text(s, x + Inches(0.25), Inches(3.18), aw - Inches(0.5), Inches(0.24), sub,
         size=11.5, color=GREY, bold=True)
    text(s, x + Inches(0.25), Inches(3.52), aw - Inches(0.5), Inches(0.6), body,
         size=12, color=INK, spacing=1.25)
    text(s, x + Inches(0.25), Inches(4.1), aw - Inches(0.5), Inches(0.26), cost,
         size=12.5, color=accent, bold=True, font=MONO)

note(s, Inches(4.75),
     [[("One knob, and we tested that it really is one knob. ", {"bold": True}),
       ("Same instructions, same analysts, same 313 meetings, same scoring — the three "
        "runs differ in nothing but the thing the model reads. We have an offline test "
        "that checks each one gets exactly what it should and nothing more, including "
        "that it is never handed a statement published after the meeting.", {})]],
     h=Inches(0.8))

punch(s, Inches(5.8),
      [[("Same brain in all three. ", {"bold": True, "color": BERK_BLUE}),
        ("So anything that moves is the input, and only the input.", {})]], size=15)

# ════════════════════════════════════════════════════════════════════════════
# 10 — Words help where the Fed speaks
# ════════════════════════════════════════════════════════════════════════════
s = slide("Words help — exactly where the Fed does the talking")

text(s, M, Inches(1.6), CONTENT_W, Inches(0.3),
     "HOW OFTEN IT CALLED THE DIRECTION RIGHT  ·  same model, only the reading changes",
     size=11, color=GREY, bold=True)

LO, HI = 0.35, 0.92
TX, TW = M + Inches(2.5), Inches(7.4)


def px(v):
    return TX + TW * ((v - LO) / (HI - LO))


def dot(sl, cx, cy, col, val, above):
    d = Inches(0.19)
    rect(sl, cx - d / 2, cy - d / 2, d, d, fill=col, shape=MSO_SHAPE.OVAL)
    text(sl, cx - Inches(0.45), cy - (Inches(0.42) if above else Inches(-0.16)),
         Inches(0.9), Inches(0.24), f"{val:.3f}", size=10, color=col, bold=True,
         align=PP_ALIGN.CENTER, font=MONO)


rows = [("Inflation", 0.533, 0.516, RED, "words made it worse"),
        ("Jobs", 0.397, 0.384, RED, "words made it worse"),
        ("Balance sheet", 0.836, 0.885, GREEN, "the Fed announces QT in words"),
        ("Term premium", 0.493, 0.543, GREEN, "a coin flip becomes a real call")]
for i, (name, a, b, col, why) in enumerate(rows):
    cy = Inches(2.35) + i * Inches(0.78)
    text(s, M, cy - Inches(0.13), Inches(2.3), Inches(0.26), name, size=13, bold=True,
         color=BERK_BLUE)
    rect(s, TX, cy - Emu(4763), TW, Emu(9525), fill=BAY_FOG)
    x1, x2 = px(a), px(b)
    rect(s, min(x1, x2), cy - Inches(0.03), abs(x2 - x1), Inches(0.06), fill=col)
    dot(s, x1, cy, ARM_VEC, a, True)
    dot(s, x2, cy, ARM_TXT, b, False)
    text(s, TX + TW + Inches(0.15), cy - Inches(0.12), Inches(2.3), Inches(0.26), why,
         size=10.5, color=col)
rect(s, px(0.5), Inches(2.12), Emu(9525), Inches(2.75), fill=PALE)
text(s, px(0.5) - Inches(1.15), Inches(1.9), Inches(1.05), Inches(0.22), "coin flip →",
     size=9.5, color=GREY, align=PP_ALIGN.RIGHT)

for i, (lab, col) in enumerate([("reading the numbers", ARM_VEC),
                                ("reading the Fed's words", ARM_TXT)]):
    x = M + Inches(2.5) + i * Inches(2.4)
    rect(s, x, Inches(5.12), Inches(0.16), Inches(0.16), fill=col, shape=MSO_SHAPE.OVAL)
    text(s, x + Inches(0.26), Inches(5.08), Inches(2.1), Inches(0.24), lab, size=11,
         color=col, bold=True)

punch(s, Inches(5.45),
      [[("The pattern isn't “text is better.” It's ", {}),
        ("“text is better where text is the source.”", {"bold": True,
                                                        "color": BERK_BLUE})],
       [("The Fed announces what it will do to its balance sheet and to long rates. It "
         "does not announce next month's CPI.", {"size": 13, "color": INK})]], size=15,
      h=Inches(0.8))

note(s, Inches(6.28),
     [[("Being honest about the evidence: ", {"bold": True}),
       ("all three runs saw the same meetings, so we can check whether a gap is bigger "
        "than luck. Only the balance sheet clears that bar — the rest point the right "
        "way on 313 meetings, but are not yet proof.", {})]], h=Inches(0.5))

# ════════════════════════════════════════════════════════════════════════════
# 11 — But words have a price
# ════════════════════════════════════════════════════════════════════════════
s = slide("But the words come with a bill")

punch(s, Inches(1.6), "Two things went wrong, and both are worse than a lower hit rate.",
      size=14, color=GREY)

# LEFT — convergence
rect(s, M, Inches(2.1), Inches(5.85), Inches(3.5), fill=LIGHT)
rect(s, M, Inches(2.1), Inches(5.85), Inches(0.07), fill=RED)
text(s, M + Inches(0.3), Inches(2.32), Inches(5.25), Inches(0.34),
     "Four specialists became one pundit", size=17, color=BERK_BLUE, bold=True)
text(s, M + Inches(0.3), Inches(2.75), Inches(5.25), Inches(0.5),
     "Feed all four analysts the same statement and they all start saying the same "
     "thing — because they are all reading the same thing.", size=12, color=INK,
     spacing=1.25)
for i, d in enumerate(["inflation", "jobs", "bal. sheet", "term prem."]):
    y = Inches(3.45) + i * Inches(0.32)
    text(s, M + Inches(0.35), y, Inches(1.2), Inches(0.24), d, size=10, color=GREY,
         align=PP_ALIGN.RIGHT)
    arrow(s, M + Inches(1.65), y + Inches(0.04), Inches(0.75), Inches(0.16), fill=PALE)
one = rect(s, M + Inches(2.55), Inches(3.5), Inches(1.5), Inches(1.05), fill=RED,
           shape=MSO_SHAPE.OVAL)
shape_text(one, "one\nvoice", 12, WHITE)
text(s, M + Inches(4.2), Inches(3.5), Inches(1.9), Inches(1.1),
     [[("0.22 → 0.34", {"size": 19, "bold": True, "color": RED, "font": MONO})],
      [("how much the four analysts now move together", {"size": 10, "color": GREY})]],
     spacing=1.2)
text(s, M + Inches(0.3), Inches(4.8), Inches(5.25), Inches(0.6),
     "The inflation analyst stops tracking inflation. Independence was the whole point "
     "of separating them — and the words quietly dissolve it.", size=12, color=RED,
     spacing=1.25)

# RIGHT — memory
xr = M + Inches(6.3)
rect(s, xr, Inches(2.1), Inches(5.83), Inches(3.5), fill=LIGHT)
rect(s, xr, Inches(2.1), Inches(5.83), Inches(0.07), fill=RED)
text(s, xr + Inches(0.3), Inches(2.32), Inches(5.23), Inches(0.34),
     "And it may be remembering, not reading", size=17, color=BERK_BLUE, bold=True)
text(s, xr + Inches(0.3), Inches(2.75), Inches(5.23), Inches(0.5),
     "The formula cannot know the future. So we count how often the model overrules it "
     "— and whether overruling it turns out to be right.", size=12, color=INK,
     spacing=1.25)
BAR_MAX = Inches(2.5)          # 353 overrules; everything scales off this
for i, (lab, n, col) in enumerate([("reading numbers", 84, ARM_VEC),
                                   ("reading the Fed's words", 353, RED)]):
    y = Inches(3.5) + i * Inches(0.62)
    text(s, xr + Inches(0.3), y, Inches(1.8), Inches(0.24), lab, size=10.5, color=GREY)
    bw_ = BAR_MAX * (n / 353)
    rect(s, xr + Inches(2.2), y, bw_, Inches(0.26), fill=col)
    text(s, xr + Inches(2.3) + bw_, y, Inches(0.8), Inches(0.26), str(n),
         size=12, color=col, bold=True, font=MONO, anchor=MSO_ANCHOR.MIDDLE)
text(s, xr + Inches(0.3), Inches(4.8), Inches(5.23), Inches(0.6),
     "4.2× as many overrules — and it gains most on exactly the two drivers where it "
     "“won”. FOMC statements are famous, dated documents. It may simply recall how "
     "yields moved next.", size=12, color=RED, spacing=1.25)

punch(s, Inches(5.8),
      [[("We cannot rule this out yet, and we say so: ", {"bold": True}),
        ("our test window ends before the model's training does, so there is no clean "
         "stretch of history it could not have seen. This is a smoke alarm, not a "
         "verdict.", {})]], size=13)

# ════════════════════════════════════════════════════════════════════════════
# 12 — Give it both
# ════════════════════════════════════════════════════════════════════════════
s = slide("Give it both: the number anchors, the words inform")

data = [
    ["", "Better than a coin flip", "Still four separate minds?", "Suspicious hindsight?"],
    ["the formula", "+0.065", "0.238   yes", "none — it can't"],
    ["numbers", "+0.065", "0.221   yes", "84 overrules"],
    ["words", "+0.082", "0.339   NO", "353 overrules"],
    ["both", "+0.086", "0.246   yes", "117 overrules"],
]
table(s, M, Inches(1.65), CONTENT_W, data, [0.22, 0.24, 0.28, 0.26],
      aligns=[L, C, C, C], size=12.5, row_h=Inches(0.42),
      fonts=[BODY, MONO, BODY, BODY],
      cell_color={(3, 1): ARM_TXT, (3, 2): RED, (3, 3): RED,
                  (4, 1): GREEN, (4, 2): GREEN, (4, 3): GREEN,
                  (2, 1): GREY, (1, 1): GREY, (2, 3): GREEN, (1, 3): GREEN,
                  (1, 2): GREEN, (2, 2): GREEN})

punch(s, Inches(4.0),
      [[("Both together is the only row that wins everywhere.", {"size": 18, "bold": True,
                                                                 "color": BERK_BLUE})],
       [("The measurement pulls each analyst back to its own driver — they are as "
         "separate as they were before the AI showed up — and the overrules fall back to "
         "a third of what the words alone provoked. The number anchors it. The story "
         "informs it.", {"size": 13.5, "color": INK})]], size=13.5, h=Inches(1.2))

rect(s, M, Inches(5.35), CONTENT_W, Inches(0.95), fill=BERK_BLUE)
text(s, M + Inches(0.35), Inches(5.5), CONTENT_W - Inches(0.7), Inches(0.7),
     [[("The input mattered more than the model.", {"size": 22, "bold": True,
                                                    "color": GOLD})],
      [("The same Claude sat behind all three columns. Everything that moved, moved "
        "because we changed what it was allowed to read.", {"size": 13, "color": WHITE})]],
     spacing=1.25, space_after=3)

text(s, M, Inches(6.45), CONTENT_W, Inches(0.3),
     [[("The whole experiment — three runs, 313 meetings each, 3,756 calls to Claude — "
        "cost ", {"size": 11.5, "color": GREY}),
       ("$4.25", {"size": 11.5, "bold": True, "color": BERK_BLUE}),
       (".  Everything we have measured so far cost $7.31.", {"size": 11.5,
                                                              "color": GREY})]])

# ════════════════════════════════════════════════════════════════════════════
# 13 — What we have not tested
# ════════════════════════════════════════════════════════════════════════════
s = slide("What we have not tested")

punch(s, Inches(1.6),
      [[("We have measured research, not money.", {"size": 20, "bold": True,
                                                   "color": RED})]], h=Inches(0.45))

items = [
    ("There is no backtest. None.",
     "No P&L, no Sharpe, on any data. The trade construction and the P&L path are "
     "written and tested — they have never been run on real history. Everything in this "
     "deck grades the analysts, not the fund."),
    ("The hindsight alarm is unresolved.",
     "Our window ends in 2024 and the model's knowledge does not. Until we score it on a "
     "stretch of history it cannot have read, “it might be remembering” stays on the "
     "table."),
    ("Three of our seven analysts add nothing.",
     "The curve, breakevens and financial-conditions analysts went in later. One is a "
     "coin flip the model overruled 3 times in 313. One is the worst analyst in the pool "
     "and loses to “assume last week repeats”."),
    ("The words experiment only covers four drivers.",
     "We ran all seven analysts once, on “both” only — so there is no numbers-vs-words "
     "comparison for the three newest. And we have never tried the Fed's minutes, only "
     "its statements."),
]
for i, (head, body) in enumerate(items):
    y = Inches(2.3) + i * Inches(1.08)
    rect(s, M, y, Inches(0.055), Inches(0.9), fill=RED if i == 0 else PALE)
    text(s, M + Inches(0.28), y, Inches(11.8), Inches(0.3), head, size=15,
         color=BERK_BLUE, bold=True)
    text(s, M + Inches(0.28), y + Inches(0.34), Inches(11.6), Inches(0.55), body,
         size=12, color=INK, spacing=1.25)

punch(s, Inches(6.6),
      [[("None of this is a reason not to show you the results — ", {}),
        ("it is the reason the results are worth anything.", {"bold": True})]], size=13)

# ════════════════════════════════════════════════════════════════════════════
# 14 — Next
# ════════════════════════════════════════════════════════════════════════════
s = slide("What we do next")

nxt = [
    ("Equity analysts, from the team", ARM_VEC,
     "Same slot, same boundary: watch one thing, say what it will do, never name a "
     "trade. If the design is right, that costs no change to the PM or the fund."),
    ("A second strategy", GOLD,
     "With one PM there is nothing to diversify and nothing to net. The whole bottom "
     "layer is written and idle until there are two books to weigh against each other."),
    ("Fix the analysts that don't work", FOUNDERS,
     "Labor is worse than a coin flip in every version we tried — the input isn't the "
     "problem there, the measurement is. Same for breakevens."),
    ("Then, finally, a backtest",
     ARM_BOTH, "Once we trust the beliefs, run the book on real history and find out "
     "whether good research survives contact with a trade."),
]
cw2 = Inches(5.85)
for i, (head, accent, body) in enumerate(nxt):
    x = M + (i % 2) * (cw2 + Inches(0.43))
    y = Inches(1.9) + (i // 2) * Inches(1.55)
    rect(s, x, y, cw2, Inches(1.3), fill=LIGHT)
    rect(s, x, y, Inches(0.055), Inches(1.3), fill=accent)
    text(s, x + Inches(0.25), y + Inches(0.16), cw2 - Inches(0.5), Inches(0.3), head,
         size=15, color=BERK_BLUE, bold=True)
    text(s, x + Inches(0.25), y + Inches(0.54), cw2 - Inches(0.5), Inches(0.65), body,
         size=12, color=INK, spacing=1.25)

rect(s, M, Inches(5.2), CONTENT_W, Inches(1.05), fill=BERK_BLUE)
text(s, M + Inches(0.35), Inches(5.38), CONTENT_W - Inches(0.7), Inches(0.8),
     [[("The one thing to take away", {"size": 11, "bold": True, "color": GOLD})],
      [("The layers let us ask which part is wrong. When we asked, the answer wasn't "
        "the model — it was what it reads.", {"size": 17, "bold": True,
                                              "color": WHITE})]], spacing=1.2,
     space_after=3)

text(s, M, Inches(6.45), CONTENT_W, Inches(0.3),
     "Appendix follows: the contracts, the transmission map, and every number behind "
     "the last six slides.", size=11.5, color=GREY)

# ════════════════════════════════════════════════════════════════════════════
# A1 — contracts
# ════════════════════════════════════════════════════════════════════════════
s = slide("The three contracts", kicker="appendix", appendix=True)
text(s, M, Inches(1.55), CONTENT_W, Inches(0.3),
     "The layers talk only through these. Any layer's method can be swapped without "
     "touching the others.", size=12, color=GREY)
cw, gap = Inches(3.85), Inches(0.29)
contracts = [
    ("DriverView", "analyst → PM", BERK_BLUE,
     [("driver", "the one force: inflation, term_premium …"),
      ("asof", "the moment the view was formed"),
      ("direction", "up / down / flat — about the DRIVER"),
      ("conviction", "0.0 → 1.0"),
      ("horizon_days", "how long it should hold"),
      ("level", "current reading — makes it scoreable"),
      ("reasoning", "why")]),
    ("StrategyTrade", "PM → fund", FOUNDERS,
     [("strategy", "e.g. macro_rates"),
      ("legs", "instrument → signed weight"),
      ("conviction", "the PM's confidence in the TRADE"),
      ("rationale", "the transmission that produced it"),
      ("risk", "what it isolates, what it hedges, net duration")]),
    ("FundAllocation", "fund → PM", BERK_BLUE,
     [("capital", "strategy → multiplier (≥ 0)"),
      ("constraints", "strategy → limits"),
      ("diagnostics", "netting, vol, breadth"),
      ("— no views —", "a control layer, not a forecasting one")]),
]
for i, (name, flow, accent, fields) in enumerate(contracts):
    x = M + i * (cw + gap)
    rect(s, x, Inches(2.0), cw, Inches(4.2), fill=WHITE, line=BAY_FOG)
    rect(s, x, Inches(2.0), cw, Inches(0.055), fill=accent)
    text(s, x + Inches(0.2), Inches(2.2), cw - Inches(0.4), Inches(0.3), name, size=16,
         color=accent, bold=True, font=MONO)
    text(s, x + Inches(0.2), Inches(2.56), cw - Inches(0.4), Inches(0.2), flow.upper(),
         size=9, color=GOLD if accent == BERK_BLUE else FOUNDERS, bold=True)
    y = Inches(2.9)
    for fname, fdesc in fields:
        text(s, x + Inches(0.2), y, cw - Inches(0.4), Inches(0.42),
             [[(fname, {"font": MONO, "size": 10, "bold": True, "color": INK}),
               ("   " + fdesc, {"size": 9.5, "color": GREY})]], spacing=1.15)
        y += Inches(0.44)

# ════════════════════════════════════════════════════════════════════════════
# A2 — the analyst + transmission map
# ════════════════════════════════════════════════════════════════════════════
s = slide("Inside an analyst, and inside the PM's map", kicker="appendix", appendix=True)

text(s, M, Inches(1.55), Inches(5.9), Inches(0.25), "THE INFLATION ANALYST, IN FULL",
     size=9.5, color=FOUNDERS, bold=True)
text(s, M, Inches(1.85), Inches(5.9), Inches(2.6),
     [[("inputs = (\"CPIAUCSL\",)   ·   horizon_days = 63",
        {"font": MONO, "size": 10, "bold": True, "color": INK})],
      [("Phase 1 — read(world) → DriverView   (deterministic, offline)",
        {"font": MONO, "size": 9.5, "bold": True, "color": BERK_BLUE})],
      [("yoy = cpi.pct_change(12);  mom = yoy[-1] − yoy[-4]",
        {"font": MONO, "size": 9, "color": INK})],
      [("direction: up / down / flat, deadband ε = 5bp", {"size": 9.5, "color": GREY})],
      [("conviction = 0.6·|mom|/1%  +  0.4·|yoy−2%|/2%", {"font": MONO, "size": 9,
                                                          "color": INK})],
      [("Phase 2 — _refine(view) → DriverView   (optional LLM)",
        {"font": MONO, "size": 9.5, "bold": True, "color": BERK_BLUE})],
      [("One call. System prompt = the persona YAML. User prompt = the reading, the "
        "FOMC statement, or both, per --input-mode. Returns JSON "
        "{direction, conviction, reasoning}. Any failure returns the Phase-1 view "
        "verbatim.", {"size": 9.5, "color": GREY})],
      [("Isolation is structural: the base class hands the analyst an AsOf gate and no "
        "other route to data. The text channel is point-in-time by each document's own "
        "release_date, not through the gate.", {"size": 9.5, "italic": True,
                                                "color": GREY})]], spacing=1.2,
     space_after=5)

xr = M + Inches(6.3)
text(s, xr, Inches(1.55), Inches(5.83), Inches(0.25),
     "THE TRANSMISSION MAP  ·  personas/macro_rates_pm.yaml", size=9.5, color=FOUNDERS,
     bold=True)
mp = [["", "front end (2y)", "long end (10y)"],
      ["inflation", "+0.50", "—"], ["labor_tightness", "+0.35", "—"],
      ["balance_sheet", "−0.15", "−0.40"], ["term_premium", "—", "+0.60"]]
table(s, xr, Inches(1.85), Inches(5.83), mp, [0.44, 0.28, 0.28], aligns=[L, C, C],
      size=10, row_h=Inches(0.3), fonts=[MONO, BODY, BODY])
text(s, xr, Inches(3.5), Inches(5.83), Inches(1.6),
     [[("slope_pressure = front − long", {"font": MONO, "size": 10.5, "bold": True,
                                          "color": BERK_BLUE}),
       ("  > 0 ⇒ flattener", {"size": 10, "color": GREY})],
      [("Legs, DV01-neutral:  +IEF (dur 7.5) · −SHY × (7.5/1.9) (dur 1.9), normalized "
        "to unit gross ⇒ net_duration ≈ 0.", {"size": 10, "color": INK})],
      [("conviction ×= (1 − 0.5 · disagreement)   — a split committee bets smaller.",
        {"font": MONO, "size": 9.5, "color": INK})],
      [("Unit-tested: |net_duration| < 0.15; hawkish ⇒ flattener, dovish ⇒ steepener; "
        "book gross ≤ leverage cap.", {"size": 9.5, "italic": True, "color": GREY})]],
     spacing=1.22, space_after=6)

# ════════════════════════════════════════════════════════════════════════════
# A3 — correctness + independence, all arms
# ════════════════════════════════════════════════════════════════════════════
s = slide("Every number: accuracy and independence", kicker="appendix", appendix=True)

text(s, M, Inches(1.5), CONTENT_W, Inches(0.25),
     "EDGE OVER A COIN FLIP  (hit_rate − 0.5), 313 meetings, 63d horizon, FRED 2019–2024",
     size=9.5, color=FOUNDERS, bold=True)
corr = [
    ["Driver", "n", "deterministic", "vector", "text", "text + vector"],
    ["inflation", "304", "+0.062", "+0.033", "+0.016", "+0.069"],
    ["labor_tightness", "242", "−0.099", "−0.103", "−0.116", "−0.074"],
    ["balance_sheet", "304", "+0.312", "+0.336", "+0.385", "+0.349"],
    ["term_premium", "300", "−0.017", "−0.007", "+0.043", "+0.000"],
    ["average edge", "", "+0.0645", "+0.0648", "+0.0820", "+0.0860"],
]
cc = {}
for r, vals in enumerate(corr[1:], start=1):
    for c in range(2, 6):
        v = float(vals[c].replace("−", "-").replace("+", ""))
        cc[(r, c)] = GREEN if v > 0.001 else (RED if v < -0.001 else GREY)
table(s, M, Inches(1.8), CONTENT_W, corr, [0.22, 0.08, 0.19, 0.17, 0.17, 0.17],
      aligns=[L, C, C, C, C, C], size=10.5, row_h=Inches(0.32),
      fonts=[MONO, BODY, BODY, BODY, BODY, BODY], cell_color=cc)

text(s, M, Inches(3.85), CONTENT_W, Inches(0.25),
     "own_corr — DOES EACH ANALYST'S VIEW STILL TRACK ITS OWN DRIVER?",
     size=9.5, color=FOUNDERS, bold=True)
ind = [
    ["Driver", "deterministic", "vector", "text", "text + vector"],
    ["inflation", "1.000", "0.965", "0.469", "0.892"],
    ["labor_tightness", "1.000", "0.902", "0.662", "0.939"],
    ["balance_sheet", "1.000", "0.881", "0.829", "0.791"],
    ["term_premium", "1.000", "0.962", "0.350", "0.927"],
    ["avg |cross-correlation|  ↓ lower = independent", "0.238", "0.221", "0.339", "0.246"],
]
ic = {(r, 3): RED for r in (1, 2, 4, 5)}
table(s, M, Inches(4.15), CONTENT_W, ind, [0.32, 0.17, 0.17, 0.17, 0.17],
      aligns=[L, C, C, C, C], size=10.5, row_h=Inches(0.32),
      fonts=[MONO, BODY, BODY, BODY, BODY], cell_color=ic)

note(s, Inches(6.1),
     [[("Faithfulness under text goes negative: ", {"bold": True}),
       ("inflation −0.171 and term_premium −0.050 — they track other drivers more than "
        "their own. Lexicon contamination under text: term_premium 0.850, inflation "
        "0.214. Source: reports/phase1_{vector,text,textvec}.md.", {})]], h=Inches(0.5))

# ════════════════════════════════════════════════════════════════════════════
# A4 — leak + significance + cost
# ════════════════════════════════════════════════════════════════════════════
s = slide("Every number: hindsight, significance, cost", kicker="appendix", appendix=True)

text(s, M, Inches(1.5), Inches(7.5), Inches(0.25),
     "OVERRIDES OF THE DETERMINISTIC CALL  ·  ovr = disagreements on gradeable dates; "
     "gain = hit-rate difference", size=9.5, color=FOUNDERS, bold=True)
lk = [
    ["Driver", "vector  ovr · gain", "text  ovr · gain", "text+vector  ovr · gain"],
    ["inflation", "26  ·  −0.029", "107  ·  −0.046", "39  ·  +0.007"],
    ["labor_tightness", "21  ·  −0.004", "87  ·  −0.017", "13  ·  +0.025"],
    ["balance_sheet", "29  ·  +0.024", "43  ·  +0.073", "51  ·  +0.037"],
    ["term_premium", "8  ·  +0.010", "116  ·  +0.060", "14  ·  +0.017"],
    ["total", "84", "353", "117"],
]
table(s, M, Inches(1.8), Inches(7.5), lk, [0.28, 0.24, 0.24, 0.24], aligns=[L, C, C, C],
      size=10, row_h=Inches(0.3), fonts=[MONO, BODY, BODY, BODY],
      cell_color={(3, 2): RED, (4, 2): RED, (5, 2): RED})

text(s, M, Inches(3.65), Inches(7.5), Inches(0.25),
     "IS THE GAP BIGGER THAN LUCK?  ·  paired McNemar, z = gain × n / √ovr",
     size=9.5, color=FOUNDERS, bold=True)
sig = [
    ["vs deterministic", "gain", "ovr", "z", ""],
    ["text × balance_sheet", "+0.073", "43", "3.4", "significant (p ≈ 0.0007)"],
    ["text × term_premium", "+0.060", "116", "1.7", "directional"],
    ["text+vec × labor", "+0.025", "13", "1.7", "directional"],
    ["text+vec × balance_sheet", "+0.037", "51", "1.6", "directional"],
    ["every other cell", "—", "—", "< 1.4", "noise"],
]
table(s, M, Inches(3.95), Inches(7.5), sig, [0.32, 0.14, 0.1, 0.1, 0.34],
      aligns=[L, R, R, R, L], size=10, row_h=Inches(0.3),
      fonts=[MONO, BODY, BODY, BODY, BODY],
      cell_color={(1, 3): GREEN, (1, 4): GREEN})

xr = M + Inches(7.9)
text(s, xr, Inches(1.5), Inches(4.23), Inches(0.25), "WHAT IT COST — AUDITED",
     size=9.5, color=FOUNDERS, bold=True)
cost = [
    ["Run", "drivers", "calls", "cost"],
    ["deterministic", "4", "0", "$0.00"],
    ["vector", "4", "1,252", "$0.94"],
    ["text", "4", "1,252", "$1.57"],
    ["text + vector", "4", "1,252", "$1.75"],
    ["text + vector", "7", "2,191", "$3.05"],
    ["total", "", "5,947", "$7.31"],
]
table(s, xr, Inches(1.8), Inches(4.23), cost, [0.36, 0.2, 0.22, 0.22],
      aligns=[L, C, R, R], size=10, row_h=Inches(0.3), fonts=[MONO, BODY, BODY, BODY])

text(s, xr, Inches(3.95), Inches(4.23), Inches(2.4),
     [[("claude-haiku-4-5  ·  0 retries  ·  17.5h wall clock  ·  4.51M input tokens",
        {"size": 10, "color": GREY})],
      [("The leak verdict the code computes", {"size": 11.5, "bold": True,
                                               "color": BERK_BLUE})],
      [("Thresholds: a gain ≤ 0.02 is no signal; override_hit needs ≥ 10 overrides to "
        "be evidence. Across 12 driver × arm cells: 1 possible leak (text × "
        "balance_sheet), 3 small-but-unconfirmed, 8 no signal.",
        {"size": 10, "color": INK})],
      [("Two caveats", {"size": 11.5, "bold": True, "color": RED})],
      [("The four older reports still print “possible training-cutoff leak” on every "
        "row — a fixed template that never read the numbers, fixed in 7b7b68c but not "
        "regenerated. And the 2019–2024 window has no post-cutoff control, so this test "
        "is inconclusive by construction.", {"size": 10, "color": INK})]], spacing=1.2,
     space_after=6)

# ════════════════════════════════════════════════════════════════════════════
# A5 — seven analysts + the offline tests
# ════════════════════════════════════════════════════════════════════════════
s = slide("The seven-analyst run, and what runs offline", kicker="appendix",
          appendix=True)

text(s, M, Inches(1.5), Inches(7.5), Inches(0.25),
     "ALL SEVEN ANALYSTS  ·  text+vector only, 2,191 calls, $3.05 — so this is NOT an A/B",
     size=9.5, color=FOUNDERS, bold=True)
seven = [
    ["Driver", "own_corr", "hit", "edge", "gain", "ovr"],
    ["inflation", "0.901", "0.559", "+0.059", "−0.003", "35"],
    ["labor_tightness", "0.935", "0.417", "−0.083", "+0.016", "13"],
    ["balance_sheet", "0.793", "0.845", "+0.345", "+0.033", "52"],
    ["term_premium", "0.930", "0.500", "+0.000", "+0.017", "14"],
    ["curve_slope   ← new", "0.973", "0.503", "+0.003", "+0.003", "3"],
    ["inflation_expectations   ← new", "0.960", "0.321", "−0.179", "0.000", "0"],
    ["financial_conditions   ← new", "0.865", "0.404", "−0.096", "+0.089", "47"],
]
table(s, M, Inches(1.8), Inches(7.5), seven, [0.34, 0.14, 0.12, 0.14, 0.14, 0.12],
      aligns=[L, C, C, C, C, C], size=10, row_h=Inches(0.3),
      fonts=[MONO, BODY, BODY, BODY, BODY, BODY],
      cell_color={(5, 3): GREY, (6, 3): RED, (7, 3): RED, (6, 2): RED, (7, 2): RED,
                  (1, 3): GREEN, (3, 3): GREEN, (2, 3): RED})
text(s, M, Inches(4.25), Inches(7.5), Inches(1.5),
     [[("None of the three adds signal.", {"size": 12, "bold": True, "color": RED})],
      [("curve_slope is a coin flip the model overrode 3 times in 313. "
        "inflation_expectations has the worst edge in the pool, zero overrides, and "
        "loses to persistence by 0.147. financial_conditions has the largest gain "
        "(+0.089) but a negative edge — it improves on a bad baseline rather than "
        "predicting.", {"size": 10.5, "color": INK})],
      [("The independence gate bit again. ", {"size": 11, "bold": True,
                                              "color": BERK_BLUE}),
       ("inflation_expectations (T10YIE) correlates 0.558 with term_premium (DGS10) — "
        "they share the real rate — and dragged term_premium's faithfulness from 0.698 "
        "to 0.442. Same lesson as the real-rates analyst (DFII10) we dropped at 0.86.",
        {"size": 10.5, "color": INK})],
      [("Caveat: the lexicon contamination of 1.000 on the three new drivers is an "
        "artifact — their mandates name other drivers, and curve_slope's own vocabulary "
        "overlaps the trade-term list. The correlation matrix shows no matching "
        "coupling.", {"size": 10, "italic": True, "color": GREY})]], spacing=1.2,
     space_after=6)

xr = M + Inches(7.9)
text(s, xr, Inches(1.5), Inches(4.23), Inches(0.25),
     "57 OFFLINE CHECKS  ·  no keys, no network, $0", size=9.5, color=FOUNDERS, bold=True)
tests = [
    ("test_layered.py", "20", "AsOf never serves data past the meeting date · each "
     "analyst reads only its own series · |net_duration| < 0.15 · hawkish ⇒ flattener, "
     "dovish ⇒ steepener · book gross ≤ leverage cap · the scorer returns 1.0 for an "
     "always-right analyst, 0.0 for an always-wrong one"),
    ("test_input_modes.py", "9", "the A/B is clean: system prompt identical across arms "
     "· the vector arm is unchanged whether or not a text source is attached · vector "
     "carries no text, text drops the numbers · point-in-time · a document dated "
     "asof+5d is never served"),
    ("test_diagnostics.py", "20", "the diagnostics harness itself, including the LLM "
     "path via a stub that manufactures overrides offline"),
    ("test_agents.py", "8", "the legacy flat ensemble degrades safely with no LLM"),
]
y = Inches(1.8)
for name, n, body in tests:
    text(s, xr, y, Inches(4.23), Inches(1.2),
         [[(name, {"font": MONO, "size": 10, "bold": True, "color": BERK_BLUE}),
           ("   " + n + " checks", {"size": 9.5, "color": GOLD, "bold": True})],
          [(body, {"size": 9.5, "color": GREY})]], spacing=1.18, space_after=3)
    y += Inches(1.28)

# ════════════════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
